#!/usr/bin/env python3
"""Build the BLAST/HMMER EnzymeX meeting deck."""

from __future__ import annotations

import argparse
import io
import zipfile
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

FONT = "Malgun Gothic"
MONO = "Consolas"

WHITE = "FFFFFF"
INK = "18212B"
CHARCOAL = "52565B"
MUTED = "667085"
LINE = "D9DEE7"
PANEL = "F5F7FA"
RED = "B71A1A"
RED_LIGHT = "FBECEC"
BLUE = "176B87"
BLUE_LIGHT = "EAF4F7"
GREEN = "167C5A"
GREEN_LIGHT = "EAF6F1"
PURPLE = "6554A4"
PURPLE_LIGHT = "F1EEFA"
AMBER = "A76500"
AMBER_LIGHT = "FFF5DF"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = INK,
    spacing: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = LINE,
    radius: bool = True,
    line_width: float = 1,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = LINE,
    width: float = 1.5,
):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def logo_blob(path: Path | None) -> bytes | None:
    if path is None or not path.exists():
        return None
    with zipfile.ZipFile(path) as archive:
        preferred = "ppt/media/image1.png"
        if preferred in archive.namelist():
            return archive.read(preferred)
    return None


def add_header(slide, title: str, number: int) -> None:
    add_rect(slide, 0, 0, SLIDE_W, 0.08, fill=RED, line=RED, radius=False, line_width=0)
    add_text(slide, title, 0.72, 0.5, 11.5, 0.62, size=30, bold=True)
    add_text(slide, str(number), 12.3, 0.62, 0.35, 0.25, size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def set_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


def add_chip(slide, text: str, x: float, y: float, w: float, *, fill: str, color: str, border: str | None = None):
    add_rect(slide, x, y, w, 0.34, fill=fill, line=border or fill, radius=True, line_width=0.7)
    add_text(slide, text, x + 0.04, y + 0.01, w - 0.08, 0.28, size=10.5, color=color,
             bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, *, accent: str = RED):
    add_rect(slide, x, y, w, 1.05, fill=WHITE, line=LINE)
    add_rect(slide, x, y, 0.09, 1.05, fill=accent, line=accent, radius=False, line_width=0)
    add_text(slide, value, x + 0.25, y + 0.16, w - 0.4, 0.44, size=26, color=accent, bold=True)
    add_text(slide, label, x + 0.25, y + 0.66, w - 0.4, 0.26, size=11.5, color=CHARCOAL)


def add_panel(slide, title: str, items: list[str], x: float, y: float, w: float, h: float,
              *, accent: str, light: str, size: float = 15):
    add_rect(slide, x, y, w, h, fill=light, line=light)
    add_text(slide, title, x + 0.3, y + 0.28, w - 0.6, 0.4, size=20, bold=True, color=accent)
    add_bullets(slide, items, x + 0.3, y + 0.85, w - 0.6, h - 1.1, size=size, spacing=11)


def add_stack_panel(slide, title: str, rows: list[tuple[str, str]], x: float, y: float,
                    w: float, *, accent: str, light: str):
    h = 0.72 + len(rows) * 0.5
    add_rect(slide, x, y, w, h, fill=WHITE, line=LINE)
    add_rect(slide, x, y, w, 0.5, fill=accent, line=accent, radius=False, line_width=0)
    add_text(slide, title, x + 0.28, y + 0.12, w - 0.56, 0.3, size=16, bold=True, color=WHITE)
    for i, (label, value) in enumerate(rows):
        row_y = y + 0.62 + i * 0.5
        if i % 2 == 0:
            add_rect(slide, x + 0.02, row_y - 0.06, w - 0.04, 0.46, fill=light, line=light,
                     radius=False, line_width=0)
        add_text(slide, label, x + 0.28, row_y, 1.55, 0.3, size=12.5, bold=True, color=CHARCOAL)
        add_text(slide, value, x + 1.95, row_y, w - 2.25, 0.34, size=12.5, color=INK)


def build_deck(output: Path, logo_deck: Path | None) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "BLAST and HMMER for EnzymeX"
    prs.core_properties.author = "Vansh Sehrawat"
    prs.core_properties.last_modified_by = "Vansh Sehrawat"
    prs.core_properties.comments = ""
    prs.core_properties.created = datetime(2026, 8, 2, 12, 0, 0)
    prs.core_properties.modified = datetime(2026, 8, 2, 12, 0, 0)
    blank = prs.slide_layouts[6]
    logo = logo_blob(logo_deck)

    # 1 ------------------------------------------------------------- title
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 0.18, SLIDE_H, fill=RED, line=RED, radius=False, line_width=0)
    if logo:
        slide.shapes.add_picture(io.BytesIO(logo), Inches(1.0), Inches(0.9), width=Inches(4.3))
    add_text(slide, "BLAST and HMMER", 1.0, 2.85, 11.0, 0.85, size=44, bold=True)
    add_chip(slide, "blastp", 1.03, 3.95, 1.12, fill=BLUE_LIGHT, color=BLUE)
    add_chip(slide, "phmmer", 2.3, 3.95, 1.22, fill=GREEN_LIGHT, color=GREEN)
    add_chip(slide, "hmmscan", 3.67, 3.95, 1.3, fill=PURPLE_LIGHT, color=PURPLE)
    add_text(slide, "Vansh Sehrawat", 1.03, 6.35, 6.0, 0.35, size=16, bold=True)
    add_text(slide, "August 4th 2026", 1.03, 6.73, 6.0, 0.32, size=14, color=MUTED)
    set_notes(slide, "Sequence-similarity evidence next to the model predictions. Not a fourth EC predictor.")

    # 2 ---------------------------------------------------------- overview
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Workflow: what happens, and where", 2)
    steps = [
        ("1", "Export", "Read the copied table once. Keep Swiss-Prot and PDB, drop bad rows, remove duplicates.",
         "PyMySQL  ->  FASTA", RED),
        ("2", "Index", "Turn that FASTA into files searchable in milliseconds, and cluster it into families.",
         "makeblastdb, MMseqs2,\nMAFFT, hmmbuild", RED),
        ("3", "Submit", "A user pastes a sequence. It is parsed and checked before anything runs.",
         "web form  ->  FASTA", BLUE),
        ("4", "Search", "All three searches run against the built files. No database is touched.",
         "blastp, phmmer, hmmscan", BLUE),
        ("5", "Show", "Attach EC, source and description to each hit, then render the tables.",
         "SQLite  ->  HTML, CSV, JSON", BLUE),
    ]
    for i, (number, title, what, tool, accent) in enumerate(steps):
        x = 0.72 + i * 2.42
        add_rect(slide, x, 1.62, 2.24, 2.72, fill=WHITE, line=LINE)
        add_rect(slide, x, 1.62, 2.24, 0.46, fill=accent, line=accent, radius=False, line_width=0)
        add_text(slide, f"{number}   {title}", x + 0.2, 1.72, 1.9, 0.28, size=14, bold=True, color=WHITE)
        add_text(slide, what, x + 0.2, 2.24, 1.86, 1.35, size=11.5, color=INK, line_spacing=1.12)
        add_line(slide, x + 0.2, 3.62, x + 2.04, 3.62, color=LINE, width=1)
        add_text(slide, tool, x + 0.2, 3.72, 1.9, 0.5, size=9.5, color=MUTED, font=MONO)
        if i < 4:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.265), Inches(2.84), Inches(0.15), Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(LINE)
            arrow.line.color.rgb = rgb(LINE)

    where = [
        (0.72, 4.66, "OFFLINE", "Runs on a schedule after each database refresh", "minutes", RED, RED_LIGHT),
        (5.56, 7.06, "ONLINE", "Runs per user request, inside the existing job queue", "seconds", BLUE, BLUE_LIGHT),
    ]
    for x, w, label, body, timing, accent, light in where:
        add_rect(slide, x, 4.56, w, 1.28, fill=light, line=light)
        add_text(slide, label, x + 0.3, 4.76, 2.0, 0.32, size=15, bold=True, color=accent)
        add_text(slide, timing, x + w - 1.5, 4.78, 1.2, 0.28, size=13, bold=True, color=accent,
                 align=PP_ALIGN.RIGHT)
        add_text(slide, body, x + 0.3, 5.2, w - 0.6, 0.42, size=13, color=INK)
    add_rect(slide, 0.72, 6.1, 11.9, 0.62, fill=PANEL, line=LINE)
    add_text(slide, "Step 1 is the only database connection in the entire flow.",
             0.72, 6.28, 11.9, 0.3, size=13.5, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    set_notes(slide, "Steps 1 and 2 are the reference build. Steps 3 to 5 are one user request. The search "
                     "code between them is a plain function with no web framework in it, so EnzymeX calls "
                     "it from its existing scheduled job rather than taking any of my web layer.")

    # 3 -------------------------------------------------------- tech stack
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Tech stack", 3)
    add_stack_panel(slide, "This project", [
        ("Language", "Python 3.11"),
        ("Web", "FastAPI + Jinja2, served by uvicorn"),
        ("Validation", "Pydantic"),
        ("Databases", "MySQL read-only for the build, SQLite for search metadata"),
        ("Search", "BLAST+ 2.16.0, HMMER 3.4"),
        ("Build tools", "MMseqs2, MAFFT"),
        ("Tests", "pytest"),
        ("Environment", "micromamba, Docker Compose or systemd"),
    ], 0.72, 1.6, 5.78, accent=RED, light=RED_LIGHT)
    add_stack_panel(slide, "EnzymeX, existing", [
        ("Language", "Python"),
        ("Web", "Pyramid, behind uWSGI and nginx"),
        ("Templates", "Jinja2"),
        ("Database", "MySQL"),
        ("Models", "ECPICK, HIT-EC, CLEAN"),
        ("Similarity", "DIAMOND, Clustal Omega"),
        ("Jobs", "scheduler with a runner table"),
        ("Data pulling", "cron, quarterly"),
    ], 6.83, 1.6, 5.79, accent=BLUE, light=BLUE_LIGHT)
    set_notes(slide, "Nothing here is exotic and nothing conflicts. The web layer I used is FastAPI only "
                     "because it was a standalone test server; the search code itself imports no web "
                     "framework, so the Pyramid side needs no rewrite. BLAST+, HMMER, MAFFT and MMseqs2 "
                     "are binaries rather than pip packages, which is why the environment is micromamba.")

    # 4 ------------------------------------------------------- two halves
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Build once, search many times", 4)
    add_panel(slide, "Offline build", [
        "The only step that opens the database",
        "Filters, deduplicates, writes FASTA",
        "Builds the BLAST index and the profiles",
        "One build ID on every artifact",
        "Minutes, on a schedule",
    ], 0.72, 1.6, 5.78, 3.5, accent=BLUE, light=BLUE_LIGHT)
    add_panel(slide, "Online search", [
        "Reads the built files only",
        "No database connection",
        "Metadata from the same build",
        "Seconds per query",
    ], 6.83, 1.6, 5.79, 3.5, accent=GREEN, light=GREEN_LIGHT)
    reasons = [
        ("No DB at request time", "cannot slow or leak the database", RED),
        ("One build ID", "hit and EC label from one snapshot", BLUE),
        ("Slow work done once", "indexing off the request path", GREEN),
    ]
    for i, (title, body, accent) in enumerate(reasons):
        x = 0.72 + i * 4.02
        add_rect(slide, x, 5.45, 3.8, 1.0, fill=WHITE, line=LINE)
        add_rect(slide, x, 5.45, 0.09, 1.0, fill=accent, line=accent, radius=False, line_width=0)
        add_text(slide, title, x + 0.28, 5.62, 3.3, 0.3, size=14, bold=True, color=accent)
        add_text(slide, body, x + 0.28, 5.97, 3.3, 0.3, size=12, color=CHARCOAL)
    set_notes(slide, "The single structural decision. Everything expensive runs offline; a request only "
                     "reads files, so a misconfigured credential is unreachable from the website.")

    # 5 ---------------------------------------------------- three searches
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Three searches", 5)
    methods = [
        ("blastp", "pairwise alignment", ["Every reference", "Identity and coverage", "Default method"], 0.72, BLUE, BLUE_LIGHT),
        ("phmmer", "query profile", ["Every reference", "Different statistics", "4-5x slower"], 4.74, GREEN, GREEN_LIGHT),
        ("hmmscan", "family profiles", ["QC-passing families only", "Finds distant relatives", "Returns a family"], 8.76, PURPLE, PURPLE_LIGHT),
    ]
    for title, sub, items, x, accent, light in methods:
        add_rect(slide, x, 1.65, 3.85, 3.55, fill=WHITE, line=LINE)
        add_rect(slide, x, 1.65, 3.85, 0.12, fill=accent, line=accent, radius=False, line_width=0)
        add_text(slide, title, x + 0.28, 1.98, 3.3, 0.42, size=24, bold=True, color=accent)
        add_text(slide, sub, x + 0.28, 2.45, 3.3, 0.3, size=12.5, color=MUTED)
        add_bullets(slide, items, x + 0.28, 2.95, 3.3, 1.9, size=14, spacing=12)
    add_rect(slide, 0.72, 5.5, 11.9, 1.0, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "Three separate tables, never merged", 0.72, 5.72, 11.9, 0.34,
             size=17, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_text(slide, "E-values are not comparable across methods", 0.72, 6.1, 11.9, 0.3,
             size=13.5, color=AMBER, align=PP_ALIGN.CENTER)
    set_notes(slide, "An E-value depends on search space size and statistical model. BLAST and phmmer "
                     "search all references, hmmscan searches the profile set. Recommendation: BLAST "
                     "always, profiles optional, phmmer off by default.")

    # 6 ---------------------------------------------------------- profiles
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Profile HMMs are built by sequence, not by EC", 6)
    add_panel(slide, "Not one profile per EC", [
        "An EC names a reaction, not a family",
        "Same reaction, unrelated folds",
        "Merging them models nothing",
    ], 0.72, 1.6, 5.78, 2.6, accent=RED, light=RED_LIGHT, size=15)
    add_panel(slide, "One profile per cluster", [
        "MMseqs2, 35% identity / 80% coverage",
        "QC gates, then MAFFT and hmmbuild",
        "EC attached after, purity reported",
    ], 6.83, 1.6, 5.79, 2.6, accent=PURPLE, light=PURPLE_LIGHT, size=15)
    add_text(slide, "One EC number, many sequence families", 0.72, 4.5, 11.9, 0.34, size=17, bold=True)
    splits = [
        ("13", "glutathione transferase", 0.72),
        ("4", "carbonic anhydrase", 4.74),
        ("5", "superoxide dismutase", 8.76),
    ]
    for value, label, x in splits:
        add_metric(slide, value, label, x, 4.95, 3.85, accent=PURPLE)
    add_rect(slide, 0.72, 6.25, 11.9, 0.55, fill=GREEN_LIGHT, line=GREEN_LIGHT)
    add_text(slide, "A human Mn-SOD query matches only the Mn/Fe families, none of the Cu/Zn ones",
             0.72, 6.4, 11.9, 0.3, size=13.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    set_notes(slide, "Superoxide dismutase has Cu/Zn, Mn/Fe and Ni enzymes under EC 1.15.1.1. A single "
                     "per-EC profile would merge unrelated folds into a model that still produces "
                     "confident-looking scores.")

    # 7 -------------------------------------------------------- validation
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Validation against the existing BLAST run", 7)
    add_rect(slide, 0.72, 1.7, 11.9, 1.85, fill=BLUE_LIGHT, line=BLUE_LIGHT)
    add_text(slide, "452 / 452", 1.1, 2.0, 4.2, 0.9, size=54, bold=True, color=BLUE)
    add_text(slide, "identical top hit", 1.15, 2.92, 4.2, 0.4, size=17, bold=True, color=BLUE)
    add_bullets(slide, [
        "Beomsu's query set, his parameters",
        "BLAST+ 2.5, evalue 1e-5, max_target_seqs 100",
    ], 5.6, 2.35, 6.8, 1.3, size=15, spacing=10)
    add_metric(slide, "450 / 452", "same top-25 subject set", 0.72, 3.95, 3.85, accent=CHARCOAL)
    add_metric(slide, "447 / 452", "same top score", 4.74, 3.95, 3.85, accent=CHARCOAL)
    add_metric(slide, "0", "provenance mismatches", 8.76, 3.95, 3.85, accent=GREEN)
    add_rect(slide, 0.72, 5.35, 11.9, 1.15, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "This is a correctness check, not an accuracy measurement",
             0.72, 5.58, 11.9, 0.34, size=17, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_text(slide, "EC agreement was measured separately, on a narrow leakage-controlled slice",
             0.72, 5.98, 11.9, 0.3, size=13.5, color=AMBER, align=PP_ALIGN.CENTER)
    set_notes(slide, "Two BLAST versions, his parameters, identical top hit on every comparable query. "
                     "Full test suite passes. On the leakage-controlled slice blastp top-hit EC agreement "
                     "was 87.7%, but that cohort is truth-selected and narrow.")

    # 8 --------------------------------------------------------- real data
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Rebuilt on real Swiss-Prot and PDB data", 8)
    add_metric(slide, "272,112", "references built", 0.72, 1.65, 3.85, accent=RED)
    add_metric(slide, "4.5 min", "offline build", 4.74, 1.65, 3.85, accent=CHARCOAL)
    add_metric(slide, "0", "search code changes", 8.76, 1.65, 3.85, accent=GREEN)
    add_text(slide, "231,577 Swiss-Prot  ·  40,535 PDB  ·  blastp 1.8 s and phmmer 7.1 s per query",
             0.72, 2.88, 11.9, 0.3, size=12.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, "What real data exposed", [
        "GFP is not an enzyme, but returns 25 hits",
        "All are PDB fusion constructs",
        "15 carry an EC, 8 different values",
        "None of them belong to GFP",
    ], 0.72, 3.35, 5.78, 3.05, accent=AMBER, light=AMBER_LIGHT, size=14.5)
    add_panel(slide, "What already catches it", [
        "Query coverage 1.00, subject coverage 0.55",
        "The gap is the fusion signature",
        "Computed from merged HSP intervals",
        "BLAST+ has no subject-coverage field",
    ], 6.83, 3.35, 5.79, 3.05, accent=BLUE, light=BLUE_LIGHT, size=14.5)
    set_notes(slide, "Public SwissProt_PDB_2022 set, a proxy for the real EnzymeX copy. Profiles were not "
                     "built on this generation, so hmmscan figures come from the smaller fixture build. "
                     "Without subject coverage the page would report a fusion partner's EC at near-perfect identity.")

    # 9 -------------------------------------------------------- what's next
    slide = prs.slides.add_slide(blank)
    add_header(slide, "What I need, and what comes next", 9)
    steps = [
        ("1", "Access", "Read-only database copy and the EnzymeX repository", RED),
        ("2", "Verify the schema", "Primary key, storage engine, source labels, UniprotID", CHARCOAL),
        ("3", "Rebuild on real data", "Re-measure, and build the profile layer", CHARCOAL),
        ("4", "Wire into the scheduler", "Alongside ECPICK, HIT-EC and CLEAN", CHARCOAL),
    ]
    for i, (number, title, body, accent) in enumerate(steps):
        y = 1.75 + i * 1.25
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.82), Inches(y), Inches(0.52), Inches(0.52))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(accent)
        circle.line.color.rgb = rgb(accent)
        add_text(slide, number, 0.83, y + 0.09, 0.5, 0.24, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            add_line(slide, 1.08, y + 0.55, 1.08, y + 1.2, color=LINE, width=2)
        add_text(slide, title, 1.62, y - 0.06, 6.9, 0.32, size=18, bold=True, color=accent)
        add_text(slide, body, 1.62, y + 0.32, 6.9, 0.3, size=13.5, color=CHARCOAL)

    add_rect(slide, 9.28, 1.65, 3.34, 4.85, fill=RED_LIGHT, line=RED_LIGHT)
    add_text(slide, "To decide", 9.6, 1.95, 2.8, 0.45, size=20, bold=True, color=RED)
    decisions = [
        "DIAMOND and BLAST both on the page, or one?",
        "phmmer by default, or optional?",
    ]
    for i, decision in enumerate(decisions):
        y = 2.7 + i * 1.35
        add_rect(slide, 9.6, y, 2.72, 1.15, fill=WHITE, line=LINE)
        add_text(slide, decision, 9.85, y + 0.15, 2.25, 0.9, size=13.5, color=INK,
                 bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Access is the blocker", 9.6, 5.6, 2.8, 0.32, size=15, bold=True, color=RED)
    set_notes(slide, "Ask for access first. Schema check and rebuild are about a day each once granted. "
                     "DIAMOND and blastp answer the same question at different sensitivity; my "
                     "recommendation is BLAST always, profiles optional, phmmer off by default.")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("output", type=Path, help="PPTX file to write")
    parser.add_argument("--logo-deck", type=Path, help="Existing EnzymeX PPTX containing ppt/media/image1.png")
    args = parser.parse_args()
    build_deck(args.output.resolve(), args.logo_deck.resolve() if args.logo_deck else None)
    print(args.output.resolve())


if __name__ == "__main__":
    main()
